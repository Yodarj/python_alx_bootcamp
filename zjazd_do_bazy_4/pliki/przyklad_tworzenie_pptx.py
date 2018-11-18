from pptx import Presentation
import csv

with open("dane/train.csv") as csvfile:
    data = csv.DictReader(csvfile, delimiter=',')
    wiek_m = 0
    surv_m = 0
    ilosc_m = 0
    ilosc_cal_m = 0
    wiek_k = 0
    ilosc_k = 0
    surv_k = 0
    ilosc_cal_k = 0
    for row in data:
        if row['Sex'] == 'male':
            ilosc_cal_m += 1
        elif row['Sex'] == 'female':
            ilosc_cal_k += 1

        if row['Survived'] == '1' and row['Sex'] == 'male':
            surv_m+=1
        elif row['Survived'] == '1' and row['Sex'] == 'female':
            surv_k+=1

        if row['Age'] == '' or row['Sex']=='':
            continue
        elif row['Sex'] == 'male':
            ilosc_m+= 1
            wiek_m += float(row['Age'])
        elif row['Sex'] == 'female':
            ilosc_k+= 1
            wiek_k += float(row['Age'])


    sredni_wiek_m = round(wiek_m/ilosc_m,2)
    sredni_wiek_k = round(wiek_k/ilosc_k,2)
    ile_m_przezylo = round(surv_m/(ilosc_cal_m),4)*100
    ile_k_przezylo = round(surv_k/(ilosc_cal_k),4)*100





prs = Presentation()

slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(slide_layout)

shapes = slide.shapes

title_shape = shapes.title

body_shape = shapes.placeholders[1]

title_shape.text = "Jakiś tekst"

tf = body_shape.text_frame
tf.text = "Statystyki"

p = tf.add_paragraph()
p.text = "Średni wiek kobiet:"
p.level = 1

p = tf.add_paragraph()
p.text = f"{sredni_wiek_k} lat"
p.level = 2

p = tf.add_paragraph()
p.text = "Ile kobiet przeżyło:"
p.level = 1

p = tf.add_paragraph()
p.text = f"{ile_k_przezylo}%"
p.level = 2

p = tf.add_paragraph()
p.text = "Średni wiek mężczyzn:"
p.level = 1

p = tf.add_paragraph()
p.text = f"{sredni_wiek_m} lat"
p.level = 2

p = tf.add_paragraph()
p.text = "Ile mężczyzn przeżyło:"
p.level = 1

p = tf.add_paragraph()
p.text = f"{ile_m_przezylo}%"
p.level = 2


prs.save("raport.pptx")